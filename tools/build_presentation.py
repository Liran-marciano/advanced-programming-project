"""
Generates PRESENTATION.pptx for the Ex 6 demo video.

Run from the repo root with: `python tools/build_presentation.py`
Produces `PRESENTATION.pptx` at the repo root with six slides:
title, background, design, live demo, beyond-the-minimum, lessons learned.
Speaker notes for each slide carry the script to read aloud while recording.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Colour palette -- matches the dark theme of the running app
# ---------------------------------------------------------------------------
BG_DARK        = RGBColor(0x0F, 0x17, 0x2A)   # slate 950
BG_PANEL       = RGBColor(0x1E, 0x29, 0x3B)   # slate 800
BG_ACCENT      = RGBColor(0x31, 0x2E, 0x81)   # indigo 900
TEXT_PRIMARY   = RGBColor(0xF1, 0xF5, 0xF9)
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)
ACCENT_CYAN    = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_PINK    = RGBColor(0xEC, 0x48, 0x99)
ACCENT_GOLD    = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_GREEN   = RGBColor(0x22, 0xC5, 0x5E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_blank_slide(prs):
    """A slide with no auto-added text placeholders, so we can lay it out by hand."""
    layout = prs.slide_layouts[6]  # 'Blank'
    return prs.slides.add_slide(layout)


def fill_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_background(slide, rgb=BG_DARK):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill_solid(rect, rgb)


def add_text(slide, left, top, width, height, text, *,
             font_size=20, bold=False, color=TEXT_PRIMARY,
             font='Calibri', align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                font_size=18, color=TEXT_PRIMARY,
                font='Calibri', bullet_color=ACCENT_CYAN,
                line_spacing=1.4):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        bullet = p.add_run()
        bullet.text = '●  '
        bullet.font.size = Pt(font_size)
        bullet.font.color.rgb = bullet_color
        bullet.font.name = font
        body = p.add_run()
        body.text = item
        body.font.size = Pt(font_size)
        body.font.color.rgb = color
        body.font.name = font
    return tb


def add_card(slide, left, top, width, height, fill=BG_PANEL):
    """A rounded rectangle that we use as a flat colour card."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_solid(card, fill)
    return card


def add_accent_bar(slide, left, top, width, height, rgb=ACCENT_CYAN):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(bar, rgb)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
def slide_title(prs):
    s = add_blank_slide(prs)
    add_background(s, BG_DARK)

    # Soft accent gradient strip on the left
    add_accent_bar(s, 0, 0, Inches(0.35), SLIDE_H, ACCENT_CYAN)

    # Course label
    add_text(s, Inches(1.0), Inches(1.0), Inches(11), Inches(0.4),
             'BAR-ILAN UNIVERSITY   •   ADVANCED PROGRAMMING (89-210 / 89-211)',
             font_size=14, bold=True, color=ACCENT_CYAN)

    # Title
    add_text(s, Inches(1.0), Inches(1.7), Inches(11), Inches(1.6),
             'Computational Graph',
             font_size=72, bold=True, color=TEXT_PRIMARY)

    add_text(s, Inches(1.0), Inches(3.4), Inches(11), Inches(0.7),
             'A browser-driven publish / subscribe pipeline',
             font_size=28, color=TEXT_SECONDARY)

    # Submitter card
    card = add_card(s, Inches(1.0), Inches(5.0), Inches(7.0), Inches(1.5))
    add_accent_bar(s, Inches(1.0), Inches(5.0), Inches(0.12), Inches(1.5), ACCENT_PURPLE)
    add_text(s, Inches(1.35), Inches(5.15), Inches(6.5), Inches(0.4),
             'Submitter', font_size=12, bold=True, color=TEXT_SECONDARY)
    add_text(s, Inches(1.35), Inches(5.45), Inches(6.5), Inches(0.5),
             'Liran Marciano', font_size=22, bold=True, color=TEXT_PRIMARY)
    add_text(s, Inches(1.35), Inches(5.95), Inches(6.5), Inches(0.4),
             'liranmar88@gmail.com', font_size=14, color=ACCENT_CYAN,
             font='Consolas')

    # Repo card
    card = add_card(s, Inches(8.3), Inches(5.0), Inches(4.0), Inches(1.5))
    add_accent_bar(s, Inches(8.3), Inches(5.0), Inches(0.12), Inches(1.5), ACCENT_GOLD)
    add_text(s, Inches(8.65), Inches(5.15), Inches(3.6), Inches(0.4),
             'Repository', font_size=12, bold=True, color=TEXT_SECONDARY)
    add_text(s, Inches(8.65), Inches(5.50), Inches(3.6), Inches(1.0),
             'github.com/\nLiran-marciano/\nadvanced-programming-project',
             font_size=11, color=ACCENT_CYAN, font='Consolas')

    set_notes(s,
        "Hi, I am Liran Marciano. This is my final project for the Advanced "
        "Programming course taught by Dr. Khalastchi at Bar-Ilan. The project "
        "brings together the six exercises we built over the semester into "
        "one live, browser-driven system for running publish/subscribe "
        "computational graphs. The repository is on GitHub at the URL on "
        "the right. Let's dive in."
    )


def slide_background(prs):
    s = add_blank_slide(prs)
    add_background(s, BG_DARK)
    add_accent_bar(s, 0, 0, Inches(0.35), SLIDE_H, ACCENT_PURPLE)

    add_text(s, Inches(0.9), Inches(0.5), Inches(11), Inches(0.4),
             'BACKGROUND', font_size=14, bold=True, color=ACCENT_PURPLE)
    add_text(s, Inches(0.9), Inches(0.9), Inches(11), Inches(0.8),
             'The pub/sub computational graph',
             font_size=36, bold=True, color=TEXT_PRIMARY)

    # Left: the idea
    card = add_card(s, Inches(0.9), Inches(2.1), Inches(5.8), Inches(4.6))
    add_text(s, Inches(1.15), Inches(2.3), Inches(5.3), Inches(0.5),
             'The idea', font_size=14, bold=True, color=ACCENT_CYAN)
    add_bullets(s, Inches(1.15), Inches(2.85), Inches(5.3), Inches(3.8), [
        'Topics are named channels.',
        'Agents subscribe to topics to receive messages, and publish to topics to notify downstream agents.',
        'Composing many agents yields a directed graph that propagates values automatically.',
        'This is the Observer pattern at scale.',
    ], font_size=17, bullet_color=ACCENT_CYAN, line_spacing=1.3)

    # Right: the example
    card = add_card(s, Inches(7.0), Inches(2.1), Inches(5.4), Inches(4.6))
    add_text(s, Inches(7.25), Inches(2.3), Inches(5.0), Inches(0.5),
             'Example: chain.conf', font_size=14, bold=True, color=ACCENT_GOLD)
    add_text(s, Inches(7.25), Inches(2.85), Inches(5.0), Inches(0.5),
             'Result = (A + B + C) + 1',
             font_size=18, bold=True, color=TEXT_PRIMARY, font='Consolas')

    add_text(s, Inches(7.25), Inches(3.55), Inches(5.0), Inches(2.8),
             ('A  ─┐\n'
              '    ├─►  PlusAgent  ─►  Sum  ─┐\n'
              'B  ─┘                                              │\n'
              '                                              └─►  PlusAgent  ─►  Total  ─►  IncAgent  ─►  Result\n'
              '                                          C  ─┘\n'),
             font_size=10, color=TEXT_SECONDARY, font='Consolas')

    set_notes(s,
        "The core idea is the observer pattern at scale. Agents subscribe "
        "to topics; when a message is published, every subscriber gets a "
        "callback and may publish derived values to its own outputs. The "
        "result is a directed graph that propagates data automatically. "
        "On the right is the example we will use throughout this demo, "
        "the chain config: it computes Result equals A plus B plus C, "
        "plus one, using two PlusAgents and one IncAgent. "
        "Exercises one to five gave us the message bus, parallel agents, "
        "the graph representation, a generic configuration loader, and a "
        "custom HTTP server. Exercise six is the view layer that turns the "
        "abstract topology into something a user can play with."
    )


def slide_design(prs):
    s = add_blank_slide(prs)
    add_background(s, BG_DARK)
    add_accent_bar(s, 0, 0, Inches(0.35), SLIDE_H, ACCENT_GOLD)

    add_text(s, Inches(0.9), Inches(0.5), Inches(11), Inches(0.4),
             'DESIGN', font_size=14, bold=True, color=ACCENT_GOLD)
    add_text(s, Inches(0.9), Inches(0.9), Inches(11), Inches(0.8),
             'Package structure & patterns',
             font_size=36, bold=True, color=TEXT_PRIMARY)

    # Left: package layout
    card = add_card(s, Inches(0.9), Inches(2.1), Inches(6.0), Inches(4.7))
    add_text(s, Inches(1.15), Inches(2.3), Inches(5.5), Inches(0.5),
             'project_biu/', font_size=14, bold=True, color=ACCENT_CYAN,
             font='Consolas')

    packages = [
        ('graph/',    'Message, Topic, ParallelAgent'),
        ('configs/',  'Graph, GenericConfig, agents'),
        ('server/',   'Homemade HTTP server (Ex 5)'),
        ('servlets/', 'TopicDisplayer, ConfLoader,\n          HtmlLoader, TopicReset'),
        ('views/',    'HtmlGraphWriter'),
        ('Main.java', 'Wires 4 servlets to 4 URIs'),
    ]
    y = Inches(2.85)
    for name, desc in packages:
        add_text(s, Inches(1.15), y, Inches(1.7), Inches(0.45),
                 name, font_size=14, bold=True, color=ACCENT_CYAN, font='Consolas')
        add_text(s, Inches(2.7), y, Inches(4.0), Inches(0.7),
                 desc, font_size=12, color=TEXT_SECONDARY)
        y += Inches(0.6)

    # Right: SOLID & patterns
    card = add_card(s, Inches(7.2), Inches(2.1), Inches(5.2), Inches(4.7))
    add_text(s, Inches(7.45), Inches(2.3), Inches(5.0), Inches(0.5),
             'Design patterns applied', font_size=14, bold=True, color=ACCENT_PURPLE)

    patterns = [
        ('Active Object', 'ParallelAgent gives every agent its own worker thread'),
        ('Singleton (Bill Pugh)', 'TopicManager via lazy inner class'),
        ('Strategy via lambda', 'BinOpAgent takes a BinaryOperator<Double>'),
        ('Template-method view', 'HtmlGraphWriter loads graph.html, injects nodes/edges'),
        ('Longest-prefix routing', 'MyHTTPServer dispatches by URI prefix'),
    ]
    y = Inches(2.85)
    for name, desc in patterns:
        add_text(s, Inches(7.45), y, Inches(4.7), Inches(0.4),
                 name, font_size=13, bold=True, color=TEXT_PRIMARY)
        add_text(s, Inches(7.45), y + Inches(0.3), Inches(4.7), Inches(0.4),
                 desc, font_size=11, color=TEXT_SECONDARY)
        y += Inches(0.75)

    set_notes(s,
        "The code is organised into six packages. The graph package holds "
        "the publish/subscribe core; the configs package describes the "
        "runtime topology; the server package is the homemade HTTP server "
        "from exercise five; the servlets package implements four request "
        "handlers; and the views package is responsible for turning a "
        "graph into HTML. "
        "Two design choices are worth pointing out. First, the view layer "
        "never builds HTML inline. It loads a static graph.html template "
        "and injects SVG fragments, so a designer can iterate on the "
        "visuals without recompiling Java. Second, every agent that comes "
        "through GenericConfig is wrapped in a ParallelAgent, so each "
        "agent runs on its own thread and the publish/subscribe layer is "
        "never blocked by slow work."
    )


def slide_demo(prs):
    s = add_blank_slide(prs)
    add_background(s, BG_DARK)
    add_accent_bar(s, 0, 0, Inches(0.35), SLIDE_H, ACCENT_GREEN)

    add_text(s, Inches(0.9), Inches(0.5), Inches(11), Inches(0.4),
             'LIVE DEMO', font_size=14, bold=True, color=ACCENT_GREEN)
    add_text(s, Inches(0.9), Inches(0.9), Inches(11), Inches(0.8),
             'Walk-through — 2 minutes',
             font_size=36, bold=True, color=TEXT_PRIMARY)

    # Five demo steps as horizontal cards
    steps = [
        ('1', 'Deploy', 'Click the simple.conf chip\nCentre fills with a 4-topic /\n2-agent graph',  ACCENT_CYAN),
        ('2', 'Publish', 'Publish A=7 and B=3\nValues cascade: C=10, D=11\nValue pills pulse-in',     ACCENT_PURPLE),
        ('3', 'Scale up', 'Click the chain.conf chip\nCanvas widens for 7 columns\nRole colours read L-to-R', ACCENT_GOLD),
        ('4', 'Reset', 'Click Reset values\nValues clear, topology stays\nA=99 alone -> Sum=99',     ACCENT_PINK),
        ('5', 'Details', 'Animated dashed edges\nLegend in the corner\nNo JS libs -- raw SVG',       ACCENT_GREEN),
    ]
    card_w = (SLIDE_W - Inches(2.0) - Inches(0.4) * 4) / 5
    x = Inches(1.0)
    y = Inches(2.3)
    for num, title, desc, accent in steps:
        add_card(s, x, y, card_w, Inches(4.5), BG_PANEL)
        add_accent_bar(s, x, y, card_w, Inches(0.18), accent)

        add_text(s, x + Inches(0.25), y + Inches(0.35), card_w - Inches(0.5), Inches(0.6),
                 num, font_size=36, bold=True, color=accent)
        add_text(s, x + Inches(0.25), y + Inches(1.05), card_w - Inches(0.5), Inches(0.5),
                 title, font_size=18, bold=True, color=TEXT_PRIMARY)
        add_text(s, x + Inches(0.25), y + Inches(1.6), card_w - Inches(0.5), Inches(2.8),
                 desc, font_size=12, color=TEXT_SECONDARY)
        x += card_w + Inches(0.4)

    # Footer hint
    add_text(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.4),
             '▶  Switch from slides to the browser at http://localhost:8080/app/index.html now',
             font_size=12, bold=True, color=ACCENT_GREEN, font='Consolas')

    set_notes(s,
        "Now I will switch to the browser. Five steps, about two minutes total.\n\n"
        "STEP 1 - I click the simple.conf chip on the left. The centre panel "
        "renders the graph: A and B are cyan input rectangles, the gold "
        "circle is PlusAgent, C is a blue intermediate topic, and D is the "
        "pink output. The right panel auto-shows the topic table.\n\n"
        "STEP 2 - I publish A equals seven, then B equals three. Notice the "
        "value pills pulse-animate above each topic the moment it updates. "
        "C is ten, D is eleven. That is two parallel agent threads firing "
        "in turn -- PlusAgent computes A plus B, then IncAgent picks up the "
        "new C and publishes C plus one.\n\n"
        "STEP 3 - Now a more complex topology. I click the chain.conf chip. "
        "Result equals A plus B plus C, plus one. The canvas resizes "
        "automatically to fit seven columns. Notice the role-based colours: "
        "cyan inputs on the left, blue intermediates in the middle, pink "
        "output on the right.\n\n"
        "STEP 4 - I click Reset values. Every topic value clears. The graph "
        "stays intact. I publish A equals ninety-nine; Sum is ninety-nine, "
        "not some number contaminated by old B, because reset also wipes "
        "every PlusAgent's internal x and y.\n\n"
        "STEP 5 - Two details I am proud of. The edges animate -- the "
        "dashed stroke offset is keyframe-animated so the graph looks "
        "alive even when nothing is happening. And the entire view is just "
        "SVG plus CSS -- no charting libraries, no canvas tag, no images."
    )


def slide_beyond(prs):
    s = add_blank_slide(prs)
    add_background(s, BG_DARK)
    add_accent_bar(s, 0, 0, Inches(0.35), SLIDE_H, ACCENT_PINK)

    add_text(s, Inches(0.9), Inches(0.5), Inches(11), Inches(0.4),
             'BEYOND THE MINIMUM', font_size=14, bold=True, color=ACCENT_PINK)
    add_text(s, Inches(0.9), Inches(0.9), Inches(11), Inches(0.8),
             'What we added past the spec',
             font_size=36, bold=True, color=TEXT_PRIMARY)

    items = [
        ('Animated, role-coloured graph',
         'Gradient nodes, drop shadows, flowing dashed edges, pulse-in value pills, colour-coded inputs / intermediates / outputs.'),
        ('Fourth servlet: TopicReset',
         'Clears every topic value and every agent\'s internal state without re-uploading the config.'),
        ('In-page example loader',
         'Two clickable chips (simple.conf, chain.conf) deploy baked-in configs in one click.'),
        ('Custom file picker',
         'Locale-proof: the native input renders "Choose file" in the browser\'s language, so we replaced it with our own.'),
        ('Dynamic canvas sizing',
         'Long pipelines (e.g. chain.conf\'s 7 columns) get a wider SVG so edges stay readable.'),
        ('Two real bugs found by the visualisation',
         'Agent name collision in Graph.createFromTopics; broken longest-path BFS fixed with Kahn\'s algorithm.'),
    ]

    cols = 2
    col_w = Inches(5.8)
    gap = Inches(0.4)
    row_h = Inches(1.35)
    y0 = Inches(2.1)
    for i, (title, desc) in enumerate(items):
        col = i % cols
        row = i // cols
        x = Inches(0.9) + col * (col_w + gap)
        y = y0 + row * row_h
        card = add_card(s, x, y, col_w, Inches(1.2))
        add_accent_bar(s, x, y, Inches(0.10), Inches(1.2),
                       ACCENT_CYAN if col == 0 else ACCENT_PURPLE)
        add_text(s, x + Inches(0.25), y + Inches(0.15), col_w - Inches(0.4), Inches(0.4),
                 title, font_size=14, bold=True, color=TEXT_PRIMARY)
        add_text(s, x + Inches(0.25), y + Inches(0.55), col_w - Inches(0.4), Inches(0.7),
                 desc, font_size=11, color=TEXT_SECONDARY)

    set_notes(s,
        "Two of these I want to highlight. First, building the view layer "
        "exposed a real design bug: every PlusAgent instance was returning "
        "the same getName, so the graph deduplicated them into a single "
        "node. Once I fixed getName to include the output topic, two "
        "distinct PlusAgents finally got two distinct circles. Second, my "
        "first layout algorithm short-circuited the longest-path "
        "relaxation; replacing the BFS with Kahn's topological sort fixed "
        "it. Both bugs were invisible to the unit tests -- only the "
        "visualisation surfaced them."
    )


def slide_lessons(prs):
    s = add_blank_slide(prs)
    add_background(s, BG_DARK)
    add_accent_bar(s, 0, 0, Inches(0.35), SLIDE_H, ACCENT_CYAN)

    add_text(s, Inches(0.9), Inches(0.5), Inches(11), Inches(0.4),
             'LESSONS LEARNED', font_size=14, bold=True, color=ACCENT_CYAN)
    add_text(s, Inches(0.9), Inches(0.9), Inches(11), Inches(0.8),
             'What I take from this course',
             font_size=36, bold=True, color=TEXT_PRIMARY)

    lessons = [
        ('SOLID is not a checklist',
         'It is the difference between adding a feature in one line and rewriting three classes.',
         ACCENT_CYAN),
        ('A good visualisation is a debugging tool',
         'The graph view surfaced two real bugs that the unit tests missed.',
         ACCENT_PURPLE),
        ('Read the PDF every time',
         'A single line ("reader.ready()") was the difference between 0 and 100 in exercise 5.',
         ACCENT_GOLD),
        ('Active Object pays off',
         'The moment one agent slows down, the whole pipeline keeps moving.',
         ACCENT_PINK),
        ('Plain stack > big frameworks',
         'Pub/sub + reflection + a homemade HTTP server is enough to ship a full product.',
         ACCENT_GREEN),
    ]

    y = Inches(2.1)
    for title, desc, accent in lessons:
        card = add_card(s, Inches(0.9), y, Inches(11.5), Inches(0.85))
        add_accent_bar(s, Inches(0.9), y, Inches(0.12), Inches(0.85), accent)
        add_text(s, Inches(1.25), y + Inches(0.10), Inches(5.0), Inches(0.4),
                 title, font_size=15, bold=True, color=TEXT_PRIMARY)
        add_text(s, Inches(6.5), y + Inches(0.10), Inches(5.8), Inches(0.7),
                 desc, font_size=12, color=TEXT_SECONDARY,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.0)

    set_notes(s,
        "To wrap up, five things stay with me. SOLID is not a checklist; "
        "applying it changes the cost of every future change. A good "
        "visualisation is a debugging tool, not just decoration. The PDF "
        "is always the source of truth -- subtle requirements buried in "
        "the text cost or saved real points more than once. The Active "
        "Object pattern pays off the moment any single agent slows down. "
        "And a plain JDK stack is enough to ship a surprising amount of "
        "product. Thank you for watching."
    )


def slide_thankyou(prs):
    s = add_blank_slide(prs)
    add_background(s, BG_DARK)

    add_text(s, 0, Inches(2.4), SLIDE_W, Inches(1.4),
             'Thank you',
             font_size=80, bold=True, color=TEXT_PRIMARY,
             align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(3.9), SLIDE_W, Inches(0.6),
             'github.com/Liran-marciano/advanced-programming-project',
             font_size=20, color=ACCENT_CYAN, font='Consolas',
             align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(4.7), SLIDE_W, Inches(0.5),
             'Liran Marciano   •   liranmar88@gmail.com',
             font_size=16, color=TEXT_SECONDARY,
             align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(5.4), SLIDE_W, Inches(0.5),
             'Advanced Programming   •   Bar-Ilan University',
             font_size=14, color=TEXT_SECONDARY,
             align=PP_ALIGN.CENTER)

    set_notes(s,
        "End-card. Hold this slide for two or three seconds, then stop "
        "the recording. Done."
    )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_background(prs)
    slide_design(prs)
    slide_demo(prs)
    slide_beyond(prs)
    slide_lessons(prs)
    slide_thankyou(prs)

    out = Path(__file__).resolve().parents[1] / 'PRESENTATION.pptx'
    prs.save(out)
    print(f'wrote {out} ({out.stat().st_size} bytes, {len(prs.slides)} slides)')


if __name__ == '__main__':
    main()
